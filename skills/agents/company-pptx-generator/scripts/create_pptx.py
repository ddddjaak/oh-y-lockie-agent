#!/usr/bin/env python3
"""
Build PPTX from layout decisions JSON.
Pure executor — no layout decisions, no height estimation, no scaling.
Every element is rendered at exactly the coordinates and sizes specified in the JSON.

Usage: python create_pptx.py <layout.json> <output.pptx> [template.pptx]
"""
import sys
import os
import re
import json

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Constants ──────────────────────────────────────────────────────────────────
FONT_CN = "微软雅黑"
FONT_EN = "Arial"
FONT_CODE = "Cascadia Code"
HEADER_BG = "0F2B46"
HEADER_TEXT = "FFFFFF"
ROW_ALT_1 = "F1F5F9"
ROW_ALT_2 = "FFFFFF"
CODE_BG = "F1F5F9"
CODE_BORDER = "CBD5E1"
CODE_TEXT = "1E293B"
TITLE_FONT_SIZE = Pt(24)

DEFAULT_BODY_SIZE = Pt(13)
DEFAULT_CODE_SIZE = Pt(9)

SLIDE_WIDTH_EMU = 12192000
SLIDE_HEIGHT_EMU = 6858000


# ── Helpers ─────────────────────────────────────────────────────────────────────

def set_run(run: object, text: str, font_name: str = FONT_CN,
            font_size: object = DEFAULT_BODY_SIZE, bold: bool = False) -> None:
    """Set text + font props on a run. Also sets East Asian font face."""
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    # Set East Asian font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT_CN)


def add_formatted_paragraph(tf: object, text: str, font_name: str = FONT_CN,
                            font_size: object = DEFAULT_BODY_SIZE, bold: bool = False,
                            bullet: bool = False, level: int = 0) -> object:
    """Add a paragraph to a text frame, parsing **bold** markers."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.level = level
    p.space_after = Pt(3)
    p.space_before = Pt(1)

    prefix = "• " if bullet else ""

    # Split on **bold** markers
    parts = re.split(r'(\*\*.*?\*\*)', text)
    has_content = False
    for part in parts:
        if not part:
            continue
        is_bold = part.startswith("**") and part.endswith("**")
        content = part[2:-2] if is_bold else part
        if not content.strip():
            continue
        run = p.add_run()
        display = prefix + content if not has_content and prefix else content
        set_run(run, display, font_name, font_size, bold=(is_bold or bold))
        has_content = True

    # If nothing was added, add empty run
    if not has_content and prefix:
        run = p.add_run()
        set_run(run, prefix, font_name, font_size)
    return p


def parse_markdown_table(markdown: str) -> tuple:
    """Parse a markdown table string into (headers, rows) tuple."""
    lines = [l.strip() for l in markdown.strip().split('\n') if l.strip()]
    if len(lines) < 2:
        return [], []

    headers = [c.strip() for c in lines[0].split('|') if c.strip()]

    rows = []
    for line in lines[1:]:
        if '---' in line and '|' in line:
            continue
        cells = [c.strip() for c in line.split('|') if c.strip()]
        if cells:
            while len(cells) < len(headers):
                cells.append("")
            rows.append(cells[:len(headers)])

    return headers, rows


def style_cell(cell: object, text: str, font_size: object = Pt(10),
               bold: bool = False, bg_color: str = None, font_color: str = None,
               align: int = PP_ALIGN.LEFT) -> None:
    """Format a table cell with font, color, margins, and background."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)

    p = tf.paragraphs[0]
    p.alignment = align
    add_formatted_paragraph(tf, text, FONT_CN, font_size, bold=bold)
    if font_color:
        for run in p.runs:
            run.font.color.rgb = RGBColor.from_string(font_color.lstrip("#"))

    # Background fill
    if bg_color:
        tcPr = cell._tc.get_or_add_tcPr()
        for fill in tcPr.findall(qn("a:solidFill")):
            tcPr.remove(fill)
        fill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgb = etree.SubElement(fill, qn("a:srgbClr"))
        srgb.set("val", bg_color.lstrip("#"))


# ── Element renderers ──────────────────────────────────────────────────────────

def render_title(slide: object, elem: dict) -> None:
    """Render title using a template placeholder by index."""
    placeholder_idx = elem.get("placeholder_idx", 0)
    text = elem.get("text", "")
    font_size_pt = elem.get("font_size_pt")

    try:
        ph = slide.placeholders[placeholder_idx]
    except (KeyError, IndexError):
        print(f"  [create_pptx] Warning: placeholder {placeholder_idx} not "
              f"found on slide, skipping title")
        return

    ph.text = text
    for p in ph.text_frame.paragraphs:
        for r in p.runs:
            run_font_size = Pt(font_size_pt) if font_size_pt else r.font.size
            set_run(r, r.text, FONT_CN, run_font_size, bold=True)


def render_text(slide: object, elem: dict) -> None:
    """Render a free text box at exact EMU coordinates."""
    text = elem.get("text", "")
    left = Emu(elem["left_emu"])
    top = Emu(elem["top_emu"])
    width = Emu(elem["width_emu"])
    height = Emu(elem["height_emu"])
    font_size = Pt(elem.get("font_size_pt", 10))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    add_formatted_paragraph(tf, text, FONT_CN, font_size)


def render_bullets(slide: object, elem: dict) -> None:
    """Render a bullet list at exact EMU coordinates."""
    items = elem.get("items", [])
    left = Emu(elem["left_emu"])
    top = Emu(elem["top_emu"])
    width = Emu(elem["width_emu"])
    height = Emu(elem["height_emu"])
    font_size = Pt(elem.get("font_size_pt", 10))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for item in items:
        add_formatted_paragraph(tf, item, FONT_CN, font_size, bullet=True)


def render_table(slide: object, elem: dict) -> None:
    """Render a styled table from a 2D data array at exact coordinates."""
    data = elem.get("data", [])
    if not data or len(data) < 2:
        print(f"  [create_pptx] Warning: table has insufficient data "
              f"(need header + at least 1 row)")
        return

    headers = data[0]
    rows_data = data[1:]
    num_rows = len(rows_data) + 1  # +1 for header row
    num_cols = len(headers)

    left = Emu(elem["left_emu"])
    top = Emu(elem["top_emu"])
    width = Emu(elem["width_emu"])
    height = Emu(elem["height_emu"])
    font_size_pt = elem.get("font_size_pt", 10)
    header_size = Pt(font_size_pt + 2)
    cell_size = Pt(font_size_pt)

    header_bg = elem.get("header_bg", HEADER_BG)
    header_text_color = elem.get("header_text", HEADER_TEXT)
    row_alt_1 = elem.get("row_alt_1", ROW_ALT_1)
    row_alt_2 = elem.get("row_alt_2", ROW_ALT_2)

    tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    tbl = tbl_shape.table

    # Evenly distribute column widths
    col_width_emu = int(width / num_cols)
    for ci in range(num_cols):
        tbl.columns[ci].width = Emu(col_width_emu)

    # Header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        style_cell(cell, str(h), header_size, bold=True,
                   bg_color=header_bg, font_color=header_text_color,
                   align=PP_ALIGN.CENTER)

    # Data rows with alternating colors
    for ri, row in enumerate(rows_data):
        bg = row_alt_1 if ri % 2 == 0 else row_alt_2
        for ci in range(num_cols):
            cell = tbl.cell(ri + 1, ci)
            text = str(row[ci]) if ci < len(row) else ""
            style_cell(cell, text, cell_size, bg_color=bg)


def render_code(slide: object, elem: dict) -> None:
    """Render a code block with styled background at exact coordinates."""
    code = elem.get("code", "")
    left = Emu(elem["left_emu"])
    top = Emu(elem["top_emu"])
    width = Emu(elem["width_emu"])
    height = Emu(elem["height_emu"])
    font_size = Pt(elem.get("font_size_pt", 9))
    font_name = elem.get("font_name", FONT_CODE)

    lines = code.split('\n') if isinstance(code, str) else code

    # Background rectangle
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(CODE_BG.lstrip("#"))
    shape.line.color.rgb = RGBColor.from_string(CODE_BORDER.lstrip("#"))
    shape.line.width = Pt(0.5)

    # Code text in the rectangle's text frame
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.add_run()
        set_run(run, line.rstrip(), font_name, font_size)
        # Set code text color
        run.font.color.rgb = RGBColor.from_string(CODE_TEXT.lstrip("#"))


def render_mermaid_image(slide: object, elem: dict) -> None:
    """Insert a mermaid PNG image, scaled to fit within specified bounds."""
    image_path = elem.get("image_path", "")
    left = Emu(elem["left_emu"])
    top = Emu(elem["top_emu"])
    width = Emu(elem["width_emu"])
    height = Emu(elem["height_emu"])

    if not image_path or not os.path.exists(image_path):
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
        tf = txBox.text_frame
        add_formatted_paragraph(tf,
                                f"[Image not found: {os.path.basename(image_path)}]",
                                FONT_CN, DEFAULT_BODY_SIZE)
        return

    try:
        from PIL import Image
        img = Image.open(image_path)
        iw, ih = img.size  # pixels (rendered at 2x scale: 1280x720 → 2560x1440)

        # Images are rendered at 2x scale. Use logical dimensions for EMU.
        MERMAID_SCALE = 2.0
        PIXEL_TO_EMU = 914400 / 96  # ~9525 EMU per pixel at 96 DPI
        img_w_emu = int(iw / MERMAID_SCALE * PIXEL_TO_EMU)
        img_h_emu = int(ih / MERMAID_SCALE * PIXEL_TO_EMU)

        # Scale to fit within specified bounds
        max_w = int(width * 0.95)
        max_h = int(height * 0.90)
        ratio = min(1.0, max_w / img_w_emu, max_h / img_h_emu)
        pw = int(img_w_emu * ratio)
        ph = int(img_h_emu * ratio)

        # Center horizontally within bounds
        cx = left + (width - pw) // 2
        slide.shapes.add_picture(image_path, cx, top, pw, ph)
    except Exception as e:
        print(f"  [create_pptx] Warning: image insert failed: {e}")


# ── Element dispatcher ─────────────────────────────────────────────────────────

ELEMENT_RENDERERS = {
    "title": render_title,
    "text": render_text,
    "bullets": render_bullets,
    "table": render_table,
    "code": render_code,
    "mermaid_image": render_mermaid_image,
}


# ── Validation ─────────────────────────────────────────────────────────────────

def validate_slide(slide_number: str, elements: list) -> None:
    """Check if any non-placeholder element overflows slide boundaries.

    Prints warnings only; never changes coordinates.
    """
    for i, elem in enumerate(elements):
        etype = elem.get("type", "unknown")
        left = elem.get("left_emu")
        top = elem.get("top_emu")
        w = elem.get("width_emu")
        h = elem.get("height_emu")

        # Title elements use placeholders — bounds are template-defined
        if left is None or top is None or w is None or h is None:
            continue

        right = left + w
        bottom = top + h

        warnings = []
        if right > SLIDE_WIDTH_EMU:
            warnings.append(
                f"right={right} exceeds slide_width={SLIDE_WIDTH_EMU} "
                f"(overflow by {right - SLIDE_WIDTH_EMU} EMU)")
        if bottom > SLIDE_HEIGHT_EMU:
            warnings.append(
                f"bottom={bottom} exceeds slide_height={SLIDE_HEIGHT_EMU} "
                f"(overflow by {bottom - SLIDE_HEIGHT_EMU} EMU)")

        for wmsg in warnings:
            print(f"  [create_pptx] WARNING: slide {slide_number}, "
                  f"element[{i}] ({etype}): {wmsg}")


# ── Placeholder cleanup ─────────────────────────────────────────────────────

def remove_unused_placeholders(slide, used_placeholder_indices):
    """Remove template placeholder shapes that were not filled with content.

    When using template layouts, slides inherit placeholders (e.g.,
    "单击此处添加文字"). After rendering our own content, unused placeholders
    must be removed to avoid visual residue.
    """
    # Collect shape elements to remove (can't modify during iteration)
    to_remove = []
    for shape in slide.shapes:
        # PLACEHOLDER type = 14
        if shape.shape_type == 14:  # MSO_SHAPE_TYPE.PLACEHOLDER
            try:
                ph_idx = shape.placeholder_format.idx
            except Exception:
                ph_idx = None

            if ph_idx not in used_placeholder_indices:
                to_remove.append(shape._element)

    # Remove from slide XML
    sp_tree = slide.shapes._spTree
    for elem in to_remove:
        sp_tree.remove(elem)


# ── Main ───────────────────────────────────────────────────────────────────────

def add_accent_line(slide):
    """Draw a thin company-red line at the top of the slide — signature element."""
    from pptx.util import Emu
    accent = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0), Emu(0),
        Emu(SLIDE_WIDTH_EMU), Emu(45000)  # ~0.05" tall, full width
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string("C00000")
    accent.line.fill.background()  # no border
    # Send to back so it doesn't cover content
    sp_tree = slide.shapes._spTree
    sp_tree.insert(2, accent._element)  # after spPr but before content shapes


def generate(json_path: str, output_path: str, template_path: str = None) -> None:
    """Generate PPTX from layout decisions JSON.

    Reads layout_decisions.json, creates slides using exact coordinates
    and font sizes from each element spec. Performs overflow validation
    after each slide but never adjusts positions.
    """
    with open(json_path, "r", encoding="utf-8") as f:
        layout_data = json.load(f)

    slides = layout_data.get("slides", [])
    if not slides:
        print("[create_pptx] Error: no slides found in layout JSON")
        sys.exit(1)

    source = layout_data.get("source", "unknown")
    print(f"[create_pptx] Layout source: {source}, {len(slides)} slides")

    # Load or create presentation
    if template_path and os.path.exists(template_path):
        print(f"[create_pptx] Using template: {template_path}")
        prs = Presentation(template_path)
        # Remove template example slides
        xml_slides = prs.slides._sldIdLst
        while len(xml_slides) > 0:
            rId = xml_slides[0].get(qn("r:id"))
            if rId:
                prs.part.drop_rel(rId)
            xml_slides.remove(xml_slides[0])
    else:
        print("[create_pptx] No template provided, creating blank widescreen "
              "presentation")
        prs = Presentation()
        prs.slide_width = Emu(SLIDE_WIDTH_EMU)
        prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    print(f"[create_pptx] Rendering {len(slides)} slides...")

    for i, slide_spec in enumerate(slides):
        slide_number = slide_spec.get("slide_number", str(i + 1))
        layout_idx = slide_spec.get("layout_idx", 0)
        layout_name = slide_spec.get("layout_name", "")
        elements = slide_spec.get("elements", [])

        # Select layout by index from JSON
        try:
            layout = prs.slide_layouts[layout_idx]
        except (IndexError, KeyError):
            print(f"  [create_pptx] Warning: layout_idx {layout_idx} not "
                  f"available, falling back to layout 0")
            layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)
        add_accent_line(slide)  # signature red PCB-trace line
        print(f"  [create_pptx] Slide {slide_number}: layout={layout_idx} "
              f"({layout_name}), {len(elements)} elements")

        # Track which placeholder indices are used
        used_placeholders = set()

        # Render each element using its exact coordinates
        for elem in elements:
            etype = elem.get("type", "")
            ph_idx = elem.get("placeholder_idx")
            if ph_idx is not None:
                used_placeholders.add(ph_idx)

            # Enforce slide boundaries: clamp element height to prevent overflow
            top_emu = elem.get("top_emu", 0)
            height_emu = elem.get("height_emu", 0)
            if top_emu and height_emu:
                max_bottom = SLIDE_HEIGHT_EMU - 50000  # 0.05" safety margin
                if top_emu + height_emu > max_bottom:
                    clamped_h = max_bottom - top_emu
                    if clamped_h > 0:
                        elem = dict(elem)  # don't mutate original
                        elem["height_emu"] = clamped_h

            renderer = ELEMENT_RENDERERS.get(etype)
            if renderer:
                try:
                    renderer(slide, elem)
                except Exception as e:
                    print(f"  [create_pptx] ERROR rendering element "
                          f"type={etype} on slide {slide_number}: {e}")
            else:
                print(f"  [create_pptx] Warning: unknown element type "
                      f"'{etype}' on slide {slide_number}, skipping")

        # Remove unused placeholder shapes (they show "单击此处添加文字")
        remove_unused_placeholders(slide, used_placeholders)

        # Validate overflow — warn only, never change
        validate_slide(slide_number, elements)

    prs.save(output_path)
    print(f"[create_pptx] Saved: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python create_pptx.py <layout.json> <output.pptx> [template.pptx]")
        sys.exit(1)

    json_path = sys.argv[1]
    output_path = sys.argv[2]
    template_path = sys.argv[3] if len(sys.argv) > 3 else None
    generate(json_path, output_path, template_path)
