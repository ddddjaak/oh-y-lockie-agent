#!/usr/bin/env python3
"""
analyze_template.py — Template analyzer for company-pptx-generator 2.0

Dual-channel extraction:
  Channel 1: python-pptx API → slide layouts, placeholder positions/types/fonts
  Channel 2: XML direct read → theme colors (clrScheme), font scheme (fontScheme),
             format scheme (fmtScheme) from ppt/theme/theme1.xml

Usage: python analyze_template.py <template.pptx> [output.json]
"""

import sys
import os
import json
import zipfile
from lxml import etree

from pptx import Presentation
from pptx.util import Pt


# ── XML Namespaces ──────────────────────────────────────────────────────────
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def _emu_to_inch(emu):
    """Convert EMU to inches (for readability in logs)."""
    return round(emu / 914400, 2) if emu else 0


# ── Channel 1: python-pptx API ──────────────────────────────────────────────

def extract_layouts(prs):
    """Extract slide layouts and their placeholders via python-pptx API."""
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "idx": i,
            "name": layout.name,
            "placeholders": [],
        }
        for ph in layout.placeholders:
            ph_type = None
            if ph.placeholder_format:
                ph_type = str(ph.placeholder_format.type)

            ph_info = {
                "idx": ph.placeholder_format.idx if ph.placeholder_format else None,
                "type": ph_type,
                "name": ph.name,
                "left_emu": ph.left,
                "top_emu": ph.top,
                "width_emu": ph.width,
                "height_emu": ph.height,
                # Inch equivalents for human readability
                "left_in": _emu_to_inch(ph.left),
                "top_in": _emu_to_inch(ph.top),
                "width_in": _emu_to_inch(ph.width),
                "height_in": _emu_to_inch(ph.height),
            }

            # Extract font info if placeholder has text
            if ph.has_text_frame:
                fonts = []
                for p in ph.text_frame.paragraphs:
                    for run in p.runs:
                        f = run.font
                        font_info = {
                            "name": f.name,
                            "size_pt": round(f.size / 12700, 1) if f.size else None,
                            "bold": f.bold,
                            "italic": f.italic,
                            "color_rgb": None,
                        }
                        # Color may be inherited from theme (NoneColor type)
                        try:
                            if f.color and f.color.type is not None:
                                font_info["color_rgb"] = str(f.color.rgb)
                        except (AttributeError, TypeError, ValueError):
                            pass
                        fonts.append(font_info)
                if fonts:
                    ph_info["fonts"] = fonts

            layout_info["placeholders"].append(ph_info)

        layouts.append(layout_info)

    return layouts


def compute_content_area(layouts):
    """
    Compute the typical content area from the most-used content layout.
    Falls back to slide dimensions minus margins if no clear content layout.
    """
    # Look for "标题和内容" or "Title and Content" layout
    content_layout = None
    for lo in layouts:
        name = lo["name"].lower()
        if any(kw in name for kw in ["标题和内容", "title and content", "标题与内容", "内容"]):
            content_layout = lo
            break

    if not content_layout:
        # Fallback: use first layout that has both TITLE and BODY
        for lo in layouts:
            types = [ph["type"] for ph in lo["placeholders"]]
            if "TITLE (1)" in types and "BODY (2)" in types:
                content_layout = lo
                break

    if not content_layout:
        return None

    # Find title bottom and content area
    title_bottom = 0
    content_top = 0
    content_left = 0
    content_width = 0
    content_bottom = 0

    # Prefer OBJECT over BODY for main content area (OBJECT is usually larger)
    body_ph = None
    object_ph = None
    for ph in content_layout["placeholders"]:
        ptype = str(ph.get("type", ""))
        if "TITLE" in ptype:
            title_bottom = max(title_bottom, ph["top_emu"] + ph["height_emu"])
        elif "BODY" in ptype and not body_ph:
            body_ph = ph
        elif "OBJECT" in ptype and not object_ph:
            object_ph = ph

    # OBJECT is typically the main content area (larger), BODY is smaller subtitle area
    main_ph = object_ph or body_ph
    if main_ph:
        content_top = main_ph["top_emu"]
        content_left = main_ph["left_emu"]
        content_width = main_ph["width_emu"]
        content_bottom = main_ph["top_emu"] + main_ph["height_emu"]

    return {
        "left_emu": content_left,
        "top_emu": content_top or title_bottom + 200000,
        "width_emu": content_width,
        "height_emu": content_bottom - (content_top or title_bottom) if content_bottom else 4800000,
        "title_bottom_emu": title_bottom,
        "source_layout": content_layout["name"],
    }


# ── Channel 2: XML Direct Read ──────────────────────────────────────────────

def extract_theme_xml(template_path):
    """
    Extract theme colors, font scheme, and format scheme directly from
    ppt/theme/theme1.xml inside the PPTX ZIP.
    """
    result = {
        "theme_colors": {},
        "font_scheme": {},
        "format_scheme": {},
    }

    try:
        with zipfile.ZipFile(template_path, 'r') as z:
            # Find the first theme XML file
            theme_files = [n for n in z.namelist()
                          if 'theme' in n.lower() and n.endswith('.xml') and '/theme' in n.lower()]

            if not theme_files:
                print(f"[analyze_template] WARN: No theme XML found in template")
                return result

            theme_xml = z.read(theme_files[0])
            root = etree.fromstring(theme_xml)

            # ── Color Scheme ──
            for clr_scheme in root.iter(f'{{{NSMAP["a"]}}}clrScheme'):
                result["theme_colors"]["name"] = clr_scheme.get("name")
                for child in clr_scheme:
                    tag = child.tag.split('}')[-1]  # e.g., dk1, lt1, accent1
                    color_val = _extract_color(child)
                    if color_val:
                        result["theme_colors"][tag] = color_val

            # ── Font Scheme ──
            for font_scheme in root.iter(f'{{{NSMAP["a"]}}}fontScheme'):
                result["font_scheme"]["name"] = font_scheme.get("name")
                for child in font_scheme:
                    tag = child.tag.split('}')[-1]  # majorFont or minorFont
                    fonts = {}
                    for sub in child:
                        script = sub.tag.split('}')[-1]  # latin, ea, cs
                        typeface = sub.get("typeface")
                        if typeface:
                            fonts[script] = typeface
                    if fonts:
                        result["font_scheme"][tag] = fonts

            # ── Format Scheme ──
            for fmt_scheme in root.iter(f'{{{NSMAP["a"]}}}fmtScheme'):
                result["format_scheme"]["name"] = fmt_scheme.get("name")
                # Extract fill style list (first solid fill color as default)
                fills = []
                for fill_style in fmt_scheme.iter(f'{{{NSMAP["a"]}}}fillStyleLst'):
                    for fill in fill_style:
                        fill_type = fill.tag.split('}')[-1]
                        fills.append(fill_type)
                result["format_scheme"]["fill_types"] = fills

    except Exception as e:
        print(f"[analyze_template] XML channel error: {e}")
        print(f"[analyze_template] Falling back to python-pptx only")

    return result


def _extract_color(element):
    """Extract color value from an <a:dk1>, <a:lt1>, <a:accent1>, etc. element."""
    for child in element:
        tag = child.tag.split('}')[-1]
        if tag == 'srgbClr':
            return child.get('val')
        elif tag == 'sysClr':
            # System color — use lastClr as fallback
            return child.get('lastClr') or child.get('val')
    return None


# ── Main ─────────────────────────────────────────────────────────────────────

def analyze_template(template_path):
    """Full template analysis: python-pptx + XML dual channel."""
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    print(f"[analyze_template] Loading: {template_path}")

    # Channel 1: python-pptx
    prs = Presentation(template_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    print(f"[analyze_template] Slide size: {slide_width} x {slide_height} EMU "
          f"({_emu_to_inch(slide_width)}\" x {_emu_to_inch(slide_height)}\")")

    layouts = extract_layouts(prs)
    print(f"[analyze_template] Channel 1 (python-pptx): {len(layouts)} layouts extracted")

    content_area = compute_content_area(layouts)
    if content_area:
        print(f"[analyze_template] Content area: {content_area['source_layout']} "
              f"({_emu_to_inch(content_area['width_emu'])}\" x {_emu_to_inch(content_area['height_emu'])}\")")

    # Channel 2: XML direct read
    theme_data = extract_theme_xml(template_path)
    if theme_data["theme_colors"]:
        n_colors = len(theme_data["theme_colors"]) - 1  # minus "name"
        print(f"[analyze_template] Channel 2 (XML): {n_colors} theme colors, "
              f"font scheme: {list(theme_data['font_scheme'].get('majorFont', {}).values())}")

    # Merge results
    result = {
        "template_path": template_path,
        "slide_width_emu": slide_width,
        "slide_height_emu": slide_height,
        "slide_width_in": _emu_to_inch(slide_width),
        "slide_height_in": _emu_to_inch(slide_height),
        "layouts": layouts,
        "content_area": content_area,
        "theme_colors": theme_data.get("theme_colors", {}),
        "font_scheme": theme_data.get("font_scheme", {}),
        "format_scheme": theme_data.get("format_scheme", {}),
    }

    # Diagnostics
    layout_names = [lo["name"] for lo in layouts]
    print(f"[analyze_template] Layout names: {layout_names}")
    if theme_data["theme_colors"]:
        colors = {k: v for k, v in theme_data["theme_colors"].items() if k != "name"}
        print(f"[analyze_template] Theme colors: {colors}")

    return result


def main():
    if len(sys.argv) < 2:
        print(f"Usage: python {sys.argv[0]} <template.pptx> [output.json]")
        sys.exit(1)

    template_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    try:
        result = analyze_template(template_path)
    except Exception as e:
        print(f"[analyze_template] ERROR: {e}")
        sys.exit(1)

    json_str = json.dumps(result, ensure_ascii=False, indent=2)

    if output_path:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(json_str)
        print(f"[analyze_template] Saved: {output_path}")
    else:
        print(json_str)


if __name__ == '__main__':
    main()
